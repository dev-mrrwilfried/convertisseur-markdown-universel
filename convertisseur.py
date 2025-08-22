#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
convertisseur.py — Convertisseur universel de fichiers vers Markdown
- Intègre les fonctionnalités de Markitdown (PPTX) et mcp-pdf-reader (PDF sémantique)
- Compatibilité Python 3.7+
- Multithread avec progression
- OCR pour PDF/images scannés
- Structure sémantique pour PDF/PPTX
- Exploration complète de sites web
"""

from __future__ import annotations
import os
import sys
import argparse
import json
import csv
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
import re
import threading
from collections import deque
from concurrent.futures import ThreadPoolExecutor, as_completed
import multiprocessing
import time
from datetime import datetime
import shutil
import requests
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
import hashlib  # Ajouté pour le hachage des noms de fichiers longs

# Nouvelle fonction utilitaire ajoutée
def is_valid_url(url: str) -> bool:
    """Valide les URLs complexes (IPv6, caractères spéciaux)"""
    try:
        result = urlparse(url)
        if all([result.scheme, result.netloc]):
            return True
        # Vérification spéciale pour les IPv6
        if '[' in result.netloc and ']' in result.netloc:
            return True
    except ValueError:
        return False
    return False

# Thread-safe print
_print_lock = threading.Lock()

def thread_safe_print(*args, **kwargs):
    with _print_lock:
        print(*args, **kwargs)

class ProgressTracker:
    """Gestionnaire de progression thread-safe"""
    
    def __init__(self, total_files: int):
        self.total_files = total_files
        self.completed = 0
        self.succeeded = 0
        self.failed = 0
        self.current_file = ""
        self.start_time = time.time()
        self.lock = threading.Lock()
        self._last_update = 0
        
    def update(self, filename: str, success: bool, message: str = ""):
        with self.lock:
            self.completed += 1
            if success:
                self.succeeded += 1
            else:
                self.failed += 1
            self.current_file = filename
            
            # Afficher la progression toutes les secondes ou à chaque fichier si peu de fichiers
            current_time = time.time()
            if (current_time - self._last_update > 1.0) or (self.total_files <= 10) or self.completed == self.total_files:
                self._display_progress()
                self._last_update = current_time
                
    def _display_progress(self):
        elapsed = time.time() - self.start_time
        percentage = (self.completed / self.total_files) * 100
        
        # Estimation du temps restant
        if self.completed > 0:
            avg_time_per_file = elapsed / self.completed
            remaining_files = self.total_files - self.completed
            eta = avg_time_per_file * remaining_files
            eta_str = f" | ETA: {self._format_time(eta)}"
        else:
            eta_str = ""
        
        # Barre de progression visuelle
        bar_width = 30
        filled = int(bar_width * percentage / 100)
        bar = "█" * filled + "░" * (bar_width - filled)
        
        # Affichage compact sur une ligne
        status = f"\r[{bar}] {percentage:5.1f}% ({self.completed}/{self.total_files}) | ✅{self.succeeded} ❌{self.failed} | {self._format_time(elapsed)}{eta_str}"
        
        # Afficher le fichier en cours si assez de place
        if len(self.current_file) < 50:
            status += f" | {os.path.basename(self.current_file)}"
        
        print(status, end="", flush=True)
        
        # Nouvelle ligne à la fin
        if self.completed == self.total_files:
            print()
            
    def _format_time(self, seconds: float) -> str:
        """Formate le temps en format lisible"""
        if seconds < 60:
            return f"{seconds:.1f}s"
        elif seconds < 3600:
            return f"{seconds/60:.1f}m"
        else:
            return f"{seconds/3600:.1f}h"
            
    def final_summary(self):
        """Affiche le résumé final"""
        elapsed = time.time() - self.start_time
        print(f"\n📊 Conversion terminée en {self._format_time(elapsed)}")
        print(f"   ✅ Réussis: {self.succeeded}")
        print(f"   ❌ Échecs: {self.failed}")
        print(f"   📈 Vitesse: {self.total_files/elapsed:.1f} fichiers/seconde")

# ---------------------------
# Utilities
# ---------------------------
def install_requirements():
    """Installe (tentative) des dépendances souvent nécessaires."""
    deps = [
        "python-docx", "beautifulsoup4", "lxml", "PyPDF2",
        "pdf2image", "pytesseract", "Pillow", "opencv-python",
        "numpy", "openpyxl", "xlrd", "python-pptx", "striprtf",
        "odfpy", "ebooklib", "pyyaml", "pdfminer.six", "pdfplumber", "requests"
    ]
    print("Installation des dépendances via pip (peut prendre du temps)...")
    for i, d in enumerate(deps, 1):
        print(f"[{i}/{len(deps)}] Installation de {d}...")
        os.system(f"{sys.executable} -m pip install {d} -q")
    print("✅ Installation terminée. Note: Tesseract (binaire) doit être installé séparément.")

def human_readable_size(size_bytes: int) -> str:
    if size_bytes == 0:
        return "0 B"
    units = ["B", "KB", "MB", "GB", "TB"]
    i = 0
    s = float(size_bytes)
    while s >= 1024 and i < len(units) - 1:
        s /= 1024.0
        i += 1
    return f"{s:.1f} {units[i]}"

def safe_import(module_name: str):
    """Import dynamique renvoyant module ou None si absent."""
    try:
        module = __import__(module_name)
        return module
    except Exception:
        return None

# ---------------------------
# Gestionnaire d'exploration de sites
# ---------------------------

class SiteCrawler:
    def __init__(self, base_url, max_depth=2, max_pages=20, delay=1, verbose=False):
        self.base_url = base_url
        self.domain = urlparse(base_url).netloc
        self.visited = set()
        self.to_visit = deque([(base_url, 0)])
        self.max_depth = max_depth
        self.max_pages = max_pages
        self.delay = delay
        self.page_count = 0
        self.verbose = verbose
        
        # Domaines de réseaux sociaux à ignorer
        self.social_domains = [
            'facebook.com', 'twitter.com', 'linkedin.com', 
            'instagram.com', 'youtube.com', 'github.com',
            'pinterest.com', 'reddit.com', 'tumblr.com',
            'snapchat.com', 'whatsapp.com', 'tiktok.com',
            'paypal.com', 't.co', 'bit.ly', 'tinyurl.com'
        ]
        
        # Mots-clés dans les URL à ignorer (version ultra-complète)
        self.ignore_keywords = [
            # Navigation et interface
            'share', 'like', 'comment', 'follow', 'subscribe',
            'header', 'footer', 'navbar', 'sidebar', 'menu',
            'navigation', 'breadcrumb', 'skip', 'anchor',
            
            # Commentaires et interactions (français + anglais)
            'comment', 'comments', 'reply', 'replies', 'respond',
            'discussion', 'feedback', 'review', 'rating', 'vote',
            'repondre', 'commenter', 'commentaire', 'commentaires',
            'reponse', 'reponses', 'avis', 'noter', 'evaluation',
            'discuter', 'reagir', 'reaction', 'reactions', 'content',
            
            # Contenu dynamique et widgets
            'content', 'widget', 'embed', 'iframe', 'popup',
            'modal', 'overlay', 'lightbox', 'carousel',
            'slider', 'gallery', 'slideshow', 'tab', 'tabs',
            'accordion', 'dropdown', 'toggle',
            
            # Interactions sociales
            'social', 'facebook', 'twitter', 'linkedin',
            'instagram', 'youtube', 'pinterest', 'whatsapp',
            'telegram', 'messenger', 'email', 'mail',
            
            # Publicités et tracking
            'advertisement', 'ads', 'banner', 'promo', 'campaign',
            'tracking', 'analytics', 'pixel', 'beacon', 'gtm',
            'publicite', 'pub', 'banniere', 'promotion',
            
            # Authentification et compte
            'login', 'signup', 'signin', 'register', 'auth',
            'logout', 'account', 'profile', 'dashboard',
            'settings', 'preferences', 'config', 'user',
            'connexion', 'inscription', 'profil', 'compte',
            'parametres', 'preferences', 'utilisateur',
            
            # E-commerce
            'cart', 'basket', 'checkout', 'payment', 'order',
            'buy', 'purchase', 'shop', 'store', 'product',
            'panier', 'commande', 'acheter', 'boutique',
            'produit', 'paiement',
            
            # Légal et informations
            'cookie', 'privacy', 'terms', 'policy', 'legal',
            'gdpr', 'rgpd', 'mentions', 'cgv', 'cgu',
            'disclaimer', 'copyright', 'license',
            
            # Navigation utilitaire
            'search', 'sitemap', 'archive', 'tag', 'category',
            'filter', 'sort', 'pagination', 'page', 'next',
            'previous', 'first', 'last', 'more',
            'recherche', 'plan', 'archives', 'categorie',
            'filtre', 'trier', 'rechercher', 'suivant',
            'precedent', 'plus',
            
            # WordPress spécifique
            'wp-admin', 'wp-content', 'wp-includes', 'wp-login',
            'xmlrpc', 'trackback', 'pingback', 'wp-json',
            
            # Autres éléments à éviter
            'contact', 'about', 'help', 'faq', 'support',
            'download', 'print', 'pdf', 'export', 'import',
            'edit', 'delete', 'admin', 'moderator',
            'donate', 'donation', 'sponsor', 'advertising',
            
            # Formats et actions
            'action=', 'do=', 'cmd=', 'mode=', 'view=',
            'output=', 'type=', 'ajax', 'json', 'xml',
            'api', 'callback', 'jsonp', 'rss', 'atom',
            'feed'
        ]
        
        # Patterns regex pour des filtres plus avancés
        self.ignore_patterns = [
            r'#content',
            r'#comment',           # Ancres vers commentaires
            r'#respond',           # Ancres de réponse
            r'#reply-\d+',         # Réponses numérotées
            r'/comment-\d+',       # URLs de commentaires
            r'/reply-\d+',         # URLs de réponses
            r'[?&]replytocom=',    # Paramètres WordPress de réponse
            r'[?&]comment_id=',    # Paramètres d'ID de commentaire
            r'[?&]reply_id=',      # Paramètres d'ID de réponse
            r'/feed/?$',           # Flux RSS/Atom
            r'/rss/?$',            # RSS
            r'/atom/?$',           # Atom
            r'\.rss$',             # Fichiers RSS
            r'\.xml$',             # Fichiers XML (souvent des flux)
            r'\.(jpg|jpeg|png|gif|webp|svg|exe|mp3|mp4|avi|mov)$'  # Fichiers média
        ]

    def get_next_url(self):
        if not self.to_visit or self.page_count >= self.max_pages:
            return None
            
        url, depth = self.to_visit.popleft()
        if url in self.visited:
            return self.get_next_url()
            
        self.visited.add(url)
        self.page_count += 1
        return url, depth

    def _should_ignore_link(self, url, link_element=None):
        """Détermine si un lien doit être ignoré"""
        parsed = urlparse(url)
        
        # Ignorer les liens vides ou avec protocoles non HTTP
        if not url or url.startswith('javascript:') or url.startswith('mailto:') or url.startswith('tel:'):
            return True
            
        # Ignorer les réseaux sociaux
        if any(domain in parsed.netloc for domain in self.social_domains):
            return True
            
        # Vérifier les mots-clés dans l'URL complète (path + query + fragment)
        full_path = (parsed.path + '?' + parsed.query + '#' + parsed.fragment).lower()
        if any(keyword in full_path for keyword in self.ignore_keywords):
            return True
            
        # Vérifier les patterns regex
        for pattern in self.ignore_patterns:
            if re.search(pattern, url.lower()):
                return True
                
        # Analyser le texte du lien et ses attributs
        if link_element:
            # Texte du lien
            link_text = link_element.get_text().lower().strip()
            ignore_texts = [
                'commentaire', 'comment', 'répondre', 'reply', 'réponse',
                'laisser un commentaire', 'leave a comment', 'post a comment',
                'share', 'partager', 'follow', 'suivre', 'like', 'aimer',
                'subscribe', 's\'abonner', 'newsletter', 'rss', 'feed'
            ]
            if any(text in link_text for text in ignore_texts):
                return True
                
            # Classes CSS du lien
            css_classes = ' '.join(link_element.get('class', [])).lower()
            ignore_classes = [
                'comment', 'reply', 'share', 'social', 'follow',
                'subscribe', 'newsletter', 'feed', 'widget',
                'sidebar', 'footer', 'header', 'nav'
            ]
            if any(cls in css_classes for cls in ignore_classes):
                return True
                
            # ID de l'élément
            element_id = link_element.get('id', '').lower()
            if any(keyword in element_id for keyword in ['comment', 'reply', 'share', 'social']):
                return True
                
        return False

    def add_links(self, soup, current_url, current_depth):
        if current_depth >= self.max_depth:
            return
            
        links_found = 0
        links_ignored = 0
        
        for link in soup.find_all('a', href=True):
            href = link['href'].strip()
            
            # Nouveau: Validation URL avant traitement
            if not href or not href.startswith(('http', '/', '#')):
                links_ignored += 1
                continue
                
            try:
                absolute_url = urljoin(current_url, href)
                
                # Validation robuste de l'URL
                if not is_valid_url(absolute_url):
                    links_ignored += 1
                    continue
                    
                parsed = urlparse(absolute_url)
                
                links_found += 1
                
                # Ne conserver que les liens du même domaine
                if parsed.netloc != self.domain:
                    links_ignored += 1
                    continue
                    
                # Vérifier si le lien doit être ignoré
                if self._should_ignore_link(absolute_url, link):
                    links_ignored += 1
                    continue
                    
                # Normaliser l'URL (supprimer les fragments sauf s'ils sont significatifs)
                clean_url = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
                if parsed.query:
                    # Garder seulement certains paramètres utiles
                    useful_params = ['id', 'p', 'page', 'post', 'article']
                    query_params = dict(param.split('=') for param in parsed.query.split('&') if '=' in param)
                    filtered_params = {k: v for k, v in query_params.items() if k in useful_params}
                    if filtered_params:
                        clean_url += '?' + '&'.join(f"{k}={v}" for k, v in filtered_params.items())
                
                # Éviter les doublons
                if (clean_url not in self.visited and 
                    clean_url not in [url for url, _ in self.to_visit] and
                    clean_url != current_url):
                    self.to_visit.append((clean_url, current_depth + 1))
                    
            except Exception as e:
                links_ignored += 1
                if self.verbose:
                    print(f"  ⚠️ Erreur traitement lien {href}: {str(e)}")
                continue
        
        # Statistiques de filtrage (optionnel)
        if links_found > 0 and self.verbose:
            print(f"  📊 Liens: {links_found} trouvés, {links_ignored} ignorés, {links_found - links_ignored} ajoutés")

# ---------------------------
# Convertisseur
# ---------------------------
class UniversalFileConverter:
    def __init__(self,
                 ocr_force: bool = False,
                 no_compress: bool = False,
                 ocr_language: str = "fra+eng",
                 threads: Optional[int] = None,
                 verbose: bool = False,
                 semantic_mode: bool = False,
                 pptx_images: bool = False,
                 pptx_notes: bool = False):
        self.ocr_force = ocr_force
        self.no_compress = no_compress
        self.ocr_language = ocr_language
        self.threads = threads or max(1, multiprocessing.cpu_count() - 1)
        self.verbose = verbose
        self.semantic_mode = semantic_mode
        self.pptx_images = pptx_images
        self.pptx_notes = pptx_notes

        # map extension -> handler
        self.supported_formats = {
            '.txt': self.convert_txt,
            '.md': self.convert_txt,
            '.rtf': self.convert_rtf,
            '.doc': self.convert_doc,
            '.docx': self.convert_docx,
            '.xls': self.convert_xls,
            '.xlsx': self.convert_xlsx,
            '.ppt': self.convert_ppt,
            '.pptx': self.convert_pptx,
            '.odt': self.convert_odt,
            '.ods': self.convert_ods,
            '.odp': self.convert_odp,
            '.html': self.convert_html,
            '.htm': self.convert_html,
            '.xml': self.convert_xml,
            '.xhtml': self.convert_html,
            '.csv': self.convert_csv,
            '.tsv': self.convert_tsv,
            '.json': self.convert_json,
            '.yaml': self.convert_yaml,
            '.yml': self.convert_yaml,
            '.pdf': self.convert_pdf,
            '.epub': self.convert_epub,
            '.png': self.convert_image,
            '.jpg': self.convert_image,
            '.jpeg': self.convert_image,
            '.bmp': self.convert_image,
            '.tiff': self.convert_image,
            '.tif': self.convert_image,
            '.gif': self.convert_image,
            '.webp': self.convert_image,
            '.log': self.convert_log,
            '.ini': self.convert_config,
            '.cfg': self.convert_config,
            '.conf': self.convert_config,
            '.url': self.convert_webpage,
            'http': self.convert_webpage,
            'https': self.convert_webpage,
        }

    @staticmethod
    def sanitize_path_component(component: str, max_length: int = 100) -> str:
        """Nettoie et tronque les composants de chemin pour éviter les erreurs"""
        # Caractères interdits sous Windows
        forbidden_chars = r'<>:"/\|?*'
        for c in forbidden_chars:
            component = component.replace(c, '_')
        # Tronquer si nécessaire
        if len(component) > max_length:
            component = component[:max_length]
        return component

    # -----------------------
    # Conversion de sites web complets
    # -----------------------
    def convert_website(self, base_url: str, output_dir: str, max_depth: int = 2, max_pages: int = 500):
        """Convertit un site web complet en Markdown avec structure hiérarchique"""
        print(f"CONVERTISSEUR URL EN MACKDOWN")
        print(f"🚀 Début de l'exploration du site: {base_url}")
        print(f"📁 Dossier de sortie: {output_dir}")
        print(f"📊 Profondeur max: {max_depth} | Pages max: {max_pages}")
        print("-" * 60)
        
        crawler = SiteCrawler(base_url, max_depth, max_pages, verbose=self.verbose)
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        # Fichier d'index
        index_content = [
            f"# Archive du site: {base_url}",
            f"**Date de création:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"**Pages converties:** 0",
            f"**Profondeur max:** {max_depth}",
            ""
        ]
        
        page_count = 0
        start_time = time.time()
        
        while True:
            next_page = crawler.get_next_url()
            if not next_page:
                break
                
            url, depth = next_page
            page_count += 1
            
            # Créer un chemin de sortie basé sur l'URL
            parsed = urlparse(url)
            domain = self.sanitize_path_component(parsed.netloc, 100)
            path_segments = [
                self.sanitize_path_component(seg, 50) 
                for seg in parsed.path.split('/') 
                if seg
            ]
            
            if not path_segments:
                file_name = "index"
                path_segments_dirs = []
            else:
                file_name = self.sanitize_path_component(path_segments[-1], 100)
                path_segments_dirs = path_segments[:-1]
            
            # Créer la structure de dossiers
            page_dir = output_path / domain
            if path_segments_dirs:
                page_dir = page_dir / '/'.join(path_segments_dirs)
            
            # Créer le chemin complet et vérifier la longueur
            full_path = page_dir / f"{file_name}.md"
            full_path_str = str(full_path)
            
            # Tronquer le chemin si trop long (nouvelle solution robuste)
            if len(full_path_str) > 240:
                # Solution 1: Raccourcir le nom de fichier
                new_file_name = self.sanitize_path_component(file_name, 50)
                full_path = page_dir / f"{new_file_name}.md"
                full_path_str = str(full_path)
                
                # Solution 2: Utiliser un hash si toujours trop long
                if len(full_path_str) > 240:
                    file_hash = hashlib.md5(url.encode('utf-8')).hexdigest()[:12]
                    full_path = page_dir / f"{file_hash}.md"
            
            try:
                full_path.parent.mkdir(parents=True, exist_ok=True)
            except Exception as e:
                print(f"❌ Erreur création dossier: {e}")
                continue
            
            # Convertir la page
            print(f"🔍 [{page_count}] Conversion: {url} (profondeur {depth})")
            markdown, soup = self.convert_webpage(url, depth)
            
            try:
                with open(full_path, 'w', encoding='utf-8') as f:
                    f.write(markdown)
            except OSError as e:
                # Gestion spécifique des erreurs de chemin trop long
                if "too long" in str(e).lower() or "nom trop long" in str(e).lower():
                    # Solution finale: utiliser un chemin court avec hash
                    file_hash = hashlib.md5(url.encode('utf-8')).hexdigest()[:12]
                    short_path = output_path / "short" / f"{file_hash}.md"
                    short_path.parent.mkdir(parents=True, exist_ok=True)
                    
                    with open(short_path, 'w', encoding='utf-8') as f:
                        f.write(markdown)
                    print(f"  ⚠️ Chemin trop long, fichier sauvegardé sous: {short_path}")
                else:
                    print(f"❌ Erreur écriture fichier: {e}")
                continue
            except Exception as e:
                print(f"❌ Erreur écriture fichier: {e}")
                continue
            
            # Ajouter à l'index
            try:
                rel_path = os.path.relpath(full_path, output_path)
                index_content.append(f"- [{parsed.path}]({rel_path}) (profondeur {depth})")
            except Exception as e:
                print(f"⚠️ Erreur création lien relatif: {e}")
                index_content.append(f"- {url} (profondeur {depth})")
            
            # Ajouter les liens trouvés pour exploration
            crawler.add_links(soup, url, depth)
            
            # Respecter un délai entre les requêtes
            time.sleep(1)
        
        # Mettre à jour le compteur de pages dans l'index
        index_content[2] = f"**Pages converties:** {page_count}"
        
        # Écrire l'index
        try:
            with open(output_path / "INDEX.md", 'w', encoding='utf-8') as f:
                f.write("\n".join(index_content))
        except Exception as e:
            print(f"❌ Erreur écriture index: {e}")
        
        elapsed = time.time() - start_time
        print("\n" + "=" * 60)
        print(f"✅ Exploration terminée en {elapsed:.1f} secondes")
        print(f"📊 Pages converties: {page_count}")
        print(f"📁 Dossier de sortie: {output_dir}")
        print(f"📄 Fichier d'index: {output_path / 'INDEX.md'}")

    # -----------------------
    # High-level file processing avec progression
    # -----------------------
    def convert_file_with_progress(self, input_path: str, output_path: Optional[str] = None, 
                                 progress_tracker: Optional[ProgressTracker] = None) -> Tuple[bool, str]:
        """
        Convertit un fichier ou une URL vers markdown avec suivi de progression.
        Retourne (success, message)
        """
        # Détecter si c'est une URL
        is_url = input_path.startswith('http://') or input_path.startswith('https://')
        
        if is_url:
            # Traitement spécial pour les URLs
            handler = self.convert_webpage
            display_name = input_path
        else:
            # Traitement normal pour les fichiers locaux
            p = Path(input_path)
            if not p.exists():
                msg = f"Fichier introuvable: {input_path}"
                if progress_tracker:
                    progress_tracker.update(p.name, False, msg)
                return False, msg

            ext = p.suffix.lower()
            handler = self.supported_formats.get(ext)
            if not handler:
                msg = f"Format non supporté : {ext}"
                if progress_tracker:
                    progress_tracker.update(p.name, False, msg)
                return False, msg
            display_name = p.name

        start_time = time.time()
        try:
            # Appel du handler approprié
            if is_url:
                md, _ = handler(input_path, 0)  # On prend uniquement le markdown pour les URLs simples
            else:
                md = handler(input_path, output_path)
        except Exception as e:
            msg = f"Erreur conversion ({display_name}): {e}"
            if progress_tracker:
                progress_tracker.update(display_name, False, msg)
            return False, msg

        # Déterminer le chemin de sortie
        if is_url:
            # Pour les URLs, générer un nom de fichier basé sur l'URL
            parsed = urlparse(input_path)
            domain = parsed.netloc.replace('www.', '')
            path = parsed.path.replace('/', '_') if parsed.path else 'homepage'
            filename = f"{domain}_{path}"[:150] + ".md"
            out = output_path or filename
        else:
            out = output_path or str(p.with_suffix('.md'))

        try:
            out_p = Path(out)
            out_p.parent.mkdir(parents=True, exist_ok=True)
            with open(out_p, 'w', encoding='utf-8') as f:
                f.write(md)
            
            # Statistiques du fichier
            conversion_time = time.time() - start_time
            if not is_url:
                input_size = p.stat().st_size
            else:
                input_size = 0  # Taille inconnue pour les URLs
            output_size = out_p.stat().st_size
            
            msg = f"Converti: {display_name} -> {out_p.name}"
            if self.verbose and not is_url:
                msg += f" ({human_readable_size(input_size)} -> {human_readable_size(output_size)}, {conversion_time:.2f}s)"
            elif self.verbose:
                msg += f" (Taille sortie: {human_readable_size(output_size)}, {conversion_time:.2f}s)"
            
            if progress_tracker:
                progress_tracker.update(display_name, True, msg)
            return True, msg
            
        except Exception as e:
            msg = f"Erreur écriture fichier: {e}"
            if progress_tracker:
                progress_tracker.update(display_name, False, msg)
            return False, msg

    # -----------------------
    # Handlers
    # -----------------------
    def convert_webpage(self, url: str, depth: int = 0) -> Tuple[str, BeautifulSoup]:
        """Convertit une page web en Markdown et retourne le contenu + soup"""
        try:
            # Télécharger le contenu de la page
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
            response = requests.get(url, headers=headers, timeout=15)
            response.raise_for_status()
            
            # Parser le contenu HTML
            soup = BeautifulSoup(response.text, 'html.parser')
            title = soup.title.string if soup.title else urlparse(url).netloc
            
            # Nettoyer le contenu
            for element in soup(['script', 'style', 'header', 'footer', 'nav', 
                               'aside', 'form', 'iframe', 'button', 'noscript']):
                element.decompose()
            
            # Structure du contenu
            content = []
            for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 
                                        'ul', 'ol', 'table', 'blockquote', 'pre']):
                # Traitement des titres
                if element.name.startswith('h'):
                    level = int(element.name[1])
                    content.append(f"{'#' * level} {element.get_text().strip()}")
                
                # Traitement des paragraphes
                elif element.name == 'p':
                    content.append(element.get_text().strip())
                
                # Traitement des listes
                elif element.name in ['ul', 'ol']:
                    list_items = []
                    for li in element.find_all('li', recursive=False):
                        prefix = '- ' if element.name == 'ul' else '1. '
                        list_items.append(f"{prefix}{li.get_text().strip()}")
                    content.append("\n".join(list_items))
                
                # Traitement des tableaux
                elif element.name == 'table':
                    rows = element.find_all('tr')
                    table_md = []
                    for i, row in enumerate(rows):
                        cells = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                        table_md.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            table_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    content.append("\n".join(table_md))
                
                # Traitement des citations et code
                elif element.name == 'blockquote':
                    content.append(f"> {element.get_text().strip()}")
                elif element.name == 'pre':
                    content.append(f"```\n{element.get_text().strip()}\n```")
            
            # Ajout des métadonnées importantes
            metadata = []
            description = soup.find('meta', attrs={'name': 'description'})
            if description:
                metadata.append(f"**Description:** {description.get('content', '')}")
                
            keywords = soup.find('meta', attrs={'name': 'keywords'})
            if keywords:
                metadata.append(f"**Keywords:** {keywords.get('content', '')}")
            
            # Construction du Markdown final
            markdown_content = [
                f"# {title}",
                f"**URL Source:** [{url}]({url})",
                f"**Date de conversion:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
                f"**Profondeur:** {depth}",
                ""
            ]
            
            if metadata:
                markdown_content.append("\n".join(metadata) + "\n")
            
            markdown_content.append("## Contenu\n")
            markdown_content.append("\n\n".join(content))
            
            return "\n".join(markdown_content), soup
        
        except Exception as e:
            return f"# Erreur de conversion web\n\n❌ Impossible de convertir {url}: {str(e)}", BeautifulSoup("", 'html.parser')

    # -----------------------
    # Méthodes de conversion - DÉBUT
    # -----------------------
    def convert_txt(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"# {p.stem}\n\n{content}"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture TXT: {e}\n"

    def convert_rtf(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from striprtf.striprtf import rtf_to_text
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install striprtf\n"
        try:
            with open(p, 'rb') as f:
                raw = f.read().decode('latin-1', errors='ignore')
            text = rtf_to_text(raw)
            return f"# {p.stem}\n\n{text}"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture RTF: {e}\n"

    def convert_docx(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from docx import Document
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install python-docx\n"
        try:
            doc = Document(str(p))
            md = f"# {p.stem}\n\n"
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    md += "\n"
                else:
                    md += text + "\n\n"
            # tableaux simples
            for table in doc.tables:
                md += "\n"
                for i, row in enumerate(table.rows):
                    cells = [c.text.strip() for c in row.cells]
                    md += "| " + " | ".join(cells) + " |\n"
                    if i == 0:
                        md += "|" + "|".join(["---" for _ in cells]) + "|\n"
                md += "\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture DOCX: {e}\n"

    def convert_doc(self, file_path: str, output_path: Optional[str] = None) -> str:
        # Lecture basique pour .doc (ancien format) : heuristique
        p = Path(file_path)
        try:
            with open(p, 'rb') as f:
                raw = f.read()
            text = raw.decode('latin-1', errors='ignore')
            # strip non-printable
            text = re.sub(r'[^\x09\x0A\x0D\x20-\x7E\u00A0-\u017F]+', ' ', text)
            return f"# {p.stem}\n\n⚠️ Conversion basique .DOC\n\n{text}"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Impossible de lire .DOC: {e}\n"

    def convert_xlsx(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from openpyxl import load_workbook
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install openpyxl\n"
        try:
            wb = load_workbook(filename=str(p), data_only=True)
            md = f"# {p.stem}\n\n"
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                md += f"## Feuille: {sheet_name}\n\n"
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    md += "*Feuille vide*\n\n"
                    continue
                headers = [str(c) if c is not None else "" for c in rows[0]]
                md += "| " + " | ".join(headers) + " |\n"
                md += "|" + "|".join(["---" for _ in headers]) + "|\n"
                for r in rows[1:101]:
                    vals = [str(c) if c is not None else "" for c in r]
                    md += "| " + " | ".join(vals) + " |\n"
                if len(rows) > 101:
                    md += f"\n*Seules les 100 premières lignes sont affichées ({len(rows)-1} lignes total)*\n"
                md += "\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture XLSX: {e}\n"

    def convert_xls(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            import xlrd
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install xlrd\n"
        try:
            wb = xlrd.open_workbook(str(p))
            md = f"# {p.stem}\n\n"
            for name in wb.sheet_names():
                md += f"## Feuille: {name}\n\n"
                sh = wb.sheet_by_name(name)
                if sh.nrows == 0:
                    md += "*Feuille vide*\n\n"
                    continue
                headers = [str(sh.cell_value(0, c)) if sh.ncols > c else f"Col{c+1}" for c in range(sh.ncols)]
                md += "| " + " | ".join(headers) + " |\n"
                md += "|" + "|".join(["---" for _ in headers]) + "|\n"
                for r in range(1, min(sh.nrows, 101)):
                    row_vals = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
                    md += "| " + " | ".join(row_vals) + " |\n"
                md += "\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture XLS: {e}\n"

    # -----------------------
    # PPTX Handler amélioré (Markitdown)
    # -----------------------
    def convert_pptx(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from pptx import Presentation
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install python-pptx\n"
        try:
            prs = Presentation(str(p))
            md = f"# {p.stem}\n\n*Présentation - {len(prs.slides)} diapositives*\n\n"
            image_counter = 1
            
            for i, slide in enumerate(prs.slides, 1):
                # Titre de la slide
                title = f"Slide {i}"
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        if shape == slide.shapes.title or "title" in shape.name.lower():
                            title = shape.text.strip()
                            break
                
                md += f"## {title}\n\n"
                
                # Contenu des formes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text and text != title:
                            md += f"{text}\n\n"
                    
                    # Extraction d'images (option --pptx-images)
                    if self.pptx_images and shape.shape_type == 13:  # Picture
                        try:
                            img_path = self._save_pptx_image(shape, output_path, i, image_counter)
                            if img_path:
                                md += f"![Image {image_counter}]({img_path})\n\n"
                                image_counter += 1
                        except Exception as e:
                            md += f"<!-- Erreur image: {e} -->\n"
                
                # Notes de présentation (option --pptx-notes)
                if self.pptx_notes:
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                        md += "### Notes\n\n"
                        md += slide.notes_slide.notes_text_frame.text.strip() + "\n\n"
                
                md += "---\n\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture PPTX: {e}\n"

    def _save_pptx_image(self, shape, output_path, slide_num, img_num):
        """Sauvegarde les images des PPTX avec chemin relatif"""
        if not output_path:
            return None
            
        output_dir = Path(output_path).parent / "images"
        output_dir.mkdir(exist_ok=True, parents=True)
        img_path = output_dir / f"slide_{slide_num}_{img_num}_{shape.name}.png"
        
        with open(img_path, "wb") as f:
            f.write(shape.image.blob)
        
        return f"./images/{img_path.name}"

    def convert_ppt(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        return f"# {p.stem}\n\n⚠️ Format .PPT (binaire ancien) non supporté directement. Convertissez en .PPTX."

    # -----------------------
    # PDF Handler amélioré (mcp-pdf-reader + OCR)
    # -----------------------
    def convert_pdf(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        
        # Mode sémantique (inspiré de mcp-pdf-reader)
        if self.semantic_mode:
            try:
                return self._convert_pdf_semantic(str(p), output_path)
            except Exception as e:
                return f"# {p.stem}\n\n❌ Erreur mode sémantique: {e}\n"
        
        # Mode OCR standard
        return self._convert_pdf_standard(str(p), output_path)
        
    def _convert_pdf_standard(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        md = f"# {p.stem}\n\n"
        # Attempt imports
        try:
            import PyPDF2
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install PyPDF2\n"

        # Optional OCR libs
        pdf2image = safe_import('pdf2image')
        pytesseract = safe_import('pytesseract')
        can_ocr = (pdf2image is not None) and (pytesseract is not None)

        try:
            with open(p, 'rb') as fh:
                reader = PyPDF2.PdfReader(fh)
                n_pages = len(reader.pages)
                md += f"*Document PDF - {n_pages} pages*\n\n"
                any_text = False
                
                # Créer un dossier pour les images si nécessaire
                images_dir = None
                if output_path:
                    output_dir = Path(output_path).parent
                    images_dir = output_dir / f"{p.stem}_images"
                    images_dir.mkdir(parents=True, exist_ok=True)
                
                for i, page in enumerate(reader.pages, 1):
                    try:
                        text = page.extract_text() or ""
                    except Exception:
                        text = ""
                    if text and text.strip():
                        any_text = True
                        md += f"## Page {i}\n\n"
                        # Basic cleanup
                        for line in text.splitlines():
                            if line.strip():
                                md += line.rstrip() + "\n\n"
                        md += "---\n\n"
                    else:
                        # Fallback OCR for this page if possible
                        if can_ocr:
                            try:
                                images = pdf2image.convert_from_path(str(p), first_page=i, last_page=i)
                                if images:
                                    img = images[0]
                                    ocr_text = pytesseract.image_to_string(img, lang=self.ocr_language)
                                    
                                    # Sauvegarder l'image originale
                                    image_path = None
                                    if images_dir:
                                        image_path = images_dir / f"page_{i:03d}.png"
                                        img.save(image_path, "PNG")
                                    
                                    if ocr_text and ocr_text.strip():
                                        any_text = True
                                        md += f"## Page {i} (OCR)\n\n"
                                        if image_path:
                                            rel_path = os.path.relpath(image_path, Path(output_path).parent)
                                            md += f"![Page {i}]({rel_path})\n\n"
                                        md += f"{ocr_text}\n\n---\n\n"
                                        continue
                                    else:
                                        # OCR a échoué mais on sauvegarde quand même l'image
                                        if image_path:
                                            rel_path = os.path.relpath(image_path, Path(output_path).parent)
                                            md += f"## Page {i}\n\n"
                                            md += f"![Page {i}]({rel_path})\n\n"
                                            md += f"*— Texte non détecté par OCR —*\n\n---\n\n"
                                        else:
                                            md += f"## Page {i}\n\n*— Aucun texte détecté (page peut être une image) —*\n\n---\n\n"
                            except Exception as e:
                                # OCR failed for this page; continue
                                md += f"<!-- OCR erreur page {i}: {e} -->\n"
                                md += f"## Page {i}\n\n*— Erreur OCR —*\n\n---\n\n"
                        # Si pas d'OCR ou OCR non disponible, on crée un lien vers l'image si possible
                        if can_ocr and images_dir:
                            try:
                                images = pdf2image.convert_from_path(str(p), first_page=i, last_page=i)
                                if images:
                                    img = images[0]
                                    image_path = images_dir / f"page_{i:03d}.png"
                                    img.save(image_path, "PNG")
                                    rel_path = os.path.relpath(image_path, Path(output_path).parent)
                                    md += f"## Page {i}\n\n"
                                    md += f"![Page {i}]({rel_path})\n\n"
                                    md += f"*— Aucun texte détecté —*\n\n---\n\n"
                                else:
                                    md += f"## Page {i}\n\n*— Aucun texte détecté (erreur extraction image) —*\n\n---\n\n"
                            except Exception:
                                md += f"## Page {i}\n\n*— Aucun texte détecté (erreur extraction image) —*\n\n---\n\n"
                        else:
                            # Not even image extraction available
                            md += f"## Page {i}\n\n*— Aucun texte détecté (page peut être une image / encodage non standard) —*\n\n---\n\n"
                if not any_text:
                    md += "\n⚠️ Aucun texte détecté dans le PDF (ni extraction directe, ni OCR).\n"
                return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture PDF: {e}\n"

    def _convert_pdf_semantic(self, file_path: str, output_path: Optional[str] = None) -> str:
        """Mode d'extraction sémantique inspiré de mcp-pdf-reader"""
        try:
            from pdfminer.high_level import extract_pages
            from pdfminer.layout import LTTextBoxHorizontal, LTFigure, LTImage
            import pdfplumber
        except ImportError:
            return f"# {Path(file_path).stem}\n\n⚠️ Dépendances manquantes pour le mode sémantique: pip install pdfminer.six pdfplumber"
        
        p = Path(file_path)
        md = [f"# {p.stem}\n\n"]
        images_dir = None
        
        if output_path:
            output_dir = Path(output_path).parent
            images_dir = output_dir / f"{p.stem}_images"
            images_dir.mkdir(parents=True, exist_ok=True)
        
        # Première passe: extraction du texte structuré
        for page_num, page_layout in enumerate(extract_pages(file_path), 1):
            md.append(f"## Page {page_num}\n\n")
            
            # Détection des éléments
            elements = []
            for element in page_layout:
                if isinstance(element, LTTextBoxHorizontal):
                    elements.append({
                        'type': 'text',
                        'text': element.get_text().strip(),
                        'x0': element.x0,
                        'y0': element.y0,
                        'size': element.size
                    })
                elif isinstance(element, LTFigure):
                    elements.append({
                        'type': 'figure',
                        'x0': element.x0,
                        'y0': element.y0
                    })
            
            # Tri par position (haut vers bas)
            elements.sort(key=lambda e: -e['y0'])
            
            # Traitement des éléments
            for elem in elements:
                if elem['type'] == 'text' and elem['text']:
                    # Détection titre par taille de police
                    if elem['size'] > 14:
                        md.append(f"### {elem['text']}\n\n")
                    else:
                        md.append(f"{elem['text']}\n\n")
            
            # Extraction des tableaux avec pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    page = pdf.pages[page_num-1]
                    tables = page.extract_tables()
                    
                    if tables:
                        md.append("\n### Tableaux\n\n")
                        for i, table in enumerate(tables, 1):
                            md.append(f"#### Tableau {i}\n\n")
                            headers = table[0]
                            md.append("| " + " | ".join(headers) + " |\n")
                            md.append("|" + "|".join(["---"] * len(headers)) + "|\n")
                            
                            for row in table[1:]:
                                md.append("| " + " | ".join(str(cell) for cell in row) + " |\n")
                            md.append("\n")
            except Exception as e:
                md.append(f"<!-- Erreur extraction tableau: {e} -->\n")
            
            md.append("---\n\n")
        
        return "".join(md)

    # -----------------------
    # Images -> OCR optional
    # -----------------------
    def convert_image(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from PIL import Image
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install Pillow\n"
        pytesseract = safe_import('pytesseract')
        can_ocr = pytesseract is not None
        try:
            img = Image.open(p)
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur ouverture image: {e}\n"
        md = f"# {p.stem}\n\n"
        md += f"**Dimensions:** {img.size[0]}x{img.size[1]}\n\n"
        if can_ocr and (self.ocr_force or True):
            try:
                text = pytesseract.image_to_string(img, lang=self.ocr_language)
                if text and text.strip():
                    md += "## Texte détecté (OCR)\n\n"
                    md += "```\n" + text.strip() + "\n```\n\n"
                else:
                    md += "*Aucun texte lisible détecté.*\n\n"
            except Exception as e:
                md += f"<!-- OCR erreur: {e} -->\n\n"
        else:
            md += "*OCR non disponible (installez pytesseract)*\n\n"
        return md

    # -----------------------
    # Autres formats
    # -----------------------
    def convert_csv(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                sample = f.read(2048)
                f.seek(0)
                sniffer = csv.Sniffer()
                dialect = sniffer.sniff(sample) if sample else csv.excel
                f.seek(0)
                reader = csv.reader(f, dialect)
                rows = list(reader)
        except Exception:
            # fallback simple split
            try:
                with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                    rows = [line.strip().split(',') for line in f.readlines()]
            except Exception as e:
                return f"# {p.stem}\n\n❌ Erreur lecture CSV: {e}\n"
        if not rows:
            return f"# {p.stem}\n\n*Fichier CSV vide*\n"
        headers = [str(h) for h in rows[0]]
        md = f"# {p.stem}\n\n| " + " | ".join(headers) + " |\n"
        md += "|" + "|".join(["---" for _ in headers]) + "|\n"
        for row in rows[1:201]:
            row = row + [""] * (len(headers) - len(row))
            md += "| " + " | ".join(str(c) for c in row[:len(headers)]) + " |\n"
        if len(rows) > 201:
            md += f"\n*Seules les 200 premières lignes affichées ({len(rows)-1} lignes total)*\n"
        return md

    def convert_tsv(self, file_path: str, output_path: Optional[str] = None) -> str:
        # Simplifié: utilise csv.reader avec tab delimiter
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                rows = list(csv.reader(f, delimiter='\t'))
            if not rows:
                return f"# {p.stem}\n\n*Fichier TSV vide*\n"
            headers = [str(h) for h in rows[0]]
            md = f"# {p.stem}\n\n| " + " | ".join(headers) + " |\n"
            md += "|" + "|".join(["---" for _ in headers]) + "|\n"
            for row in rows[1:201]:
                row = row + [""] * (len(headers) - len(row))
                md += "| " + " | ".join(str(c) for c in row[:len(headers)]) + " |\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture TSV: {e}\n"

    def convert_json(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8') as f:
                data = json.load(f)
            pretty = json.dumps(data, indent=2, ensure_ascii=False)
            return f"# {p.stem}\n\n```json\n{pretty}\n```"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur JSON: {e}\n"

    def convert_yaml(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            import yaml
            with open(p, 'r', encoding='utf-8') as f:
                content = f.read()
            return f"# {p.stem}\n\n```yaml\n{content}\n```"
        except Exception:
            # fallback raw
            try:
                with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return f"# {p.stem}\n\n```yaml\n{content}\n```"
            except Exception as e:
                return f"# {p.stem}\n\n❌ Erreur lecture YAML: {e}\n"

    def convert_html(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from bs4 import BeautifulSoup
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install beautifulsoup4\n"
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                html = f.read()
            soup = BeautifulSoup(html, 'html.parser')
            for s in soup(["script", "style"]):
                s.decompose()
            text = soup.get_text(separator="\n")
            md = f"# {p.stem}\n\n{text.strip()}\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture HTML: {e}\n"

    def convert_xml(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            tree = ET.parse(p)
            root = tree.getroot()
            xml_str = ET.tostring(root, encoding='unicode')
            return f"# {p.stem}\n\n```xml\n{xml_str}\n```"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur XML: {e}\n"

    def convert_epub(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendances manquantes: pip install ebooklib beautifulsoup4\n"
        try:
            book = epub.read_epub(str(p))
            md = f"# {p.stem}\n\n"
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text(separator="\n").strip()
                    if text:
                        md += text + "\n\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur EPUB: {e}\n"

    def convert_log(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"# Log: {p.stem}\n\n```\n{content}\n```"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture log: {e}\n"

    def convert_config(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            with open(p, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"# Config: {p.stem}\n\n```\n{content}\n```"
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur lecture config: {e}\n"

    def convert_odt(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        try:
            from odf.opendocument import load
            from odf.text import P
        except Exception:
            return f"# {p.stem}\n\n⚠️ Dépendance manquante: pip install odfpy\n"
        try:
            doc = load(str(p))
            md = f"# {p.stem}\n\n"
            for elem in doc.getElementsByType(P):
                text = "".join([node.data for node in elem.childNodes if getattr(node, "data", None)])
                if text.strip():
                    md += text.strip() + "\n\n"
            return md
        except Exception as e:
            return f"# {p.stem}\n\n❌ Erreur ODT: {e}\n"

    def convert_ods(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        return f"# {p.stem}\n\n⚠️ Conversion ODS non implémentée (exportez en XLSX pour meilleure compatibilité)\n"

    def convert_odp(self, file_path: str, output_path: Optional[str] = None) -> str:
        p = Path(file_path)
        return f"# {p.stem}\n\n⚠️ Conversion ODP non implémentée (exportez en PPTX)\n"

    # -----------------------
    # Directory / batch processing avec progression
    # -----------------------
    def convert_directory(self, dir_path: str, output_dir: Optional[str] = None, recursive: bool = False):
        input_dir = Path(dir_path)
        if not input_dir.is_dir():
            print(f"❌ Répertoire introuvable: {dir_path}")
            return

        out_dir = Path(output_dir) if output_dir else input_dir / "markdown_output"
        out_dir.mkdir(parents=True, exist_ok=True)

        pattern = "**/*" if recursive else "*"
        files = [p for p in input_dir.glob(pattern) if p.is_file() and p.suffix.lower() in self.supported_formats]
        
        if not files:
            print("Aucun fichier supporté trouvé.")
            return

        # Analyse préliminaire
        total_size = sum(p.stat().st_size for p in files)
        print(f"🔍 Analyse: {len(files)} fichier(s) trouvé(s)")
        print(f"📦 Taille totale: {human_readable_size(total_size)}")
        print(f"🧵 Threads: {self.threads}")
        print(f"📁 Sortie: {out_dir}")
        print(f"🔄 Mode: {'Récursif' if recursive else 'Direct'}")
        print("-" * 60)

        # Initialiser le tracker de progression
        progress_tracker = ProgressTracker(len(files))

        with ThreadPoolExecutor(max_workers=self.threads) as ex:
            futures = {}
            for p in files:
                rel = p.relative_to(input_dir)
                out = out_dir / rel.with_suffix('.md')
                out.parent.mkdir(parents=True, exist_ok=True)
                futures[ex.submit(self.convert_file_with_progress, str(p), str(out), progress_tracker)] = (p, out)
            
            # Attendre la completion
            for fut in as_completed(futures):
                p, out = futures[fut]
                try:
                    ok, msg = fut.result()
                except Exception as e:
                    progress_tracker.update(p.name, False, f"Exception: {e}")

        progress_tracker.final_summary()
        print(f"📁 Fichiers de sortie disponibles dans: {out_dir}")

    def batch_convert(self, files: List[str], output_dir: Optional[str] = None):
        if not files:
            print("Aucun fichier fourni.")
            return
            
        # Normalize list, filter unsupported
        valid = []
        invalid = []
        for f in files:
            p = Path(f)
            if p.exists() and p.suffix.lower() in self.supported_formats:
                valid.append(p)
            else:
                invalid.append(f)
                
        if invalid:
            print("⚠️ Fichiers ignorés (inexistant ou format non supporté):")
            for i in invalid:
                print(f"  - {i}")
                
        if not valid:
            print("❌ Aucun fichier valide à convertir.")
            return

        # Analyse préliminaire
        total_size = sum(p.stat().st_size for p in valid)
        print(f"🔍 Analyse: {len(valid)} fichier(s) valide(s)")
        print(f"📦 Taille totale: {human_readable_size(total_size)}")
        print(f"🧵 Threads: {self.threads}")
        if output_dir:
            print(f"📁 Sortie: {output_dir}")
        print("-" * 60)

        # Initialiser le tracker de progression
        progress_tracker = ProgressTracker(len(valid))
        
        with ThreadPoolExecutor(max_workers=self.threads) as ex:
            futures = {}
            for p in valid:
                out_path = str(Path(output_dir) / f"{p.stem}.md") if output_dir else None
                futures[ex.submit(self.convert_file_with_progress, str(p), out_path, progress_tracker)] = p
                
            # Attendre la completion
            for fut in as_completed(futures):
                p = futures[fut]
                try:
                    ok, msg = fut.result()
                except Exception as e:
                    progress_tracker.update(p.name, False, f"Exception: {e}")

        progress_tracker.final_summary()

# ---------------------------
# CLI
# ---------------------------
def parse_args():
    parser = argparse.ArgumentParser(description="Convertisseur universel de fichiers vers Markdown (multithread avec progression).")
    parser.add_argument("input", nargs="*", help="Fichier(s) ou dossier(s) à convertir.")
    parser.add_argument("-o", "--output", help="Fichier ou dossier de sortie (si multiple fichiers -> dossier).")
    parser.add_argument("-d", "--directory", action="store_true", help="Traiter le paramètre input comme répertoire (prendre input[0]).")
    parser.add_argument("-r", "--recursive", action="store_true", help="Recherche récursive dans les répertoires.")
    parser.add_argument("--batch", action="store_true", help="Mode batch pour plusieurs fichiers listés.")
    parser.add_argument("--install-deps", action="store_true", help="Installe les dépendances Python utiles (pip).")
    parser.add_argument("--ocr-only", action="store_true", help="Force l'OCR pour les images (active ocr_force).")
    parser.add_argument("--no-compress", action="store_true", help="Désactive la compression des images.")
    parser.add_argument("--ocr-lang", default="fra+eng", help="Langues pour l'OCR (ex: fra+eng, eng, spa).")
    parser.add_argument("--threads", type=int, default=None, help="Nombre de threads (par défaut CPU-1).")
    parser.add_argument("-v", "--verbose", action="store_true", help="Affichage détaillé (tailles fichiers, temps de conversion).")
    parser.add_argument("--no-progress", action="store_true", help="Désactive la barre de progression (mode silencieux).")
    
    # Nouvelles options
    parser.add_argument("--semantic", action="store_true", help="Mode d'extraction sémantique pour PDF (structure hiérarchique).")
    parser.add_argument("--pptx-images", action="store_true", help="Extraire les images des fichiers PPTX.")
    parser.add_argument("--pptx-notes", action="store_true", help="Extraire les notes des présentations PPTX.")
    
    # Option pour l'exploration de sites
    parser.add_argument("--website", action="store_true", help="Exploration complète d'un site web")
    parser.add_argument("--depth", type=int, default=2, help="Profondeur d'exploration pour les sites web")
    parser.add_argument("--max-pages", type=int, default=50, help="Nombre maximum de pages à convertir")
    
    return parser.parse_args()

def main():
    args = parse_args()

    if args.install_deps:
        install_requirements()
        return

    conv = UniversalFileConverter(
        ocr_force=args.ocr_only,
        no_compress=args.no_compress,
        ocr_language=args.ocr_lang,
        threads=args.threads,
        verbose=args.verbose,
        semantic_mode=args.semantic,
        pptx_images=args.pptx_images,
        pptx_notes=args.pptx_notes
    )

    # Mode site web
    if args.website and args.input:
        conv.convert_website(
            base_url=args.input[0],
            output_dir=args.output or "site_archive",
            max_depth=args.depth,
            max_pages=args.max_pages
        )
        return

    if not args.input:
        print("❌ Aucun fichier ou dossier fourni. Utilisez --help pour l'aide.")
        return

    # Afficher les informations de démarrage
    if not args.no_progress:
        print(f"🚀 Convertisseur Document - {datetime.now().strftime('%H:%M:%S')}")
        print(f"💻 CPU: {multiprocessing.cpu_count()} cœurs disponibles")
        if args.semantic:
            print("🔍 Mode sémantique activé pour les PDF")
        if args.pptx_images or args.pptx_notes:
            pptx_modes = []
            if args.pptx_images: pptx_modes.append("images")
            if args.pptx_notes: pptx_modes.append("notes")
            print(f"📊 Mode PPTX amélioré: {', '.join(pptx_modes)}")

    # Directory mode
    if args.directory:
        input_path = args.input[0]
        conv.convert_directory(input_path, output_dir=args.output, recursive=args.recursive)
        return

    # Batch mode for multiple files
    if args.batch or len(args.input) > 1:
        # if output provided and is a directory ensure it exists
        if args.output:
            out_p = Path(args.output)
            if out_p.exists() and not out_p.is_dir():
                print("❌ L'option -o/--output doit être un dossier lorsque plusieurs fichiers sont fournis.")
                return
            out_p.mkdir(parents=True, exist_ok=True)
        conv.batch_convert(args.input, output_dir=args.output)
        return

    # Single file mode
    input_path = args.input[0]
    out = args.output
    
    if not args.no_progress:
        p = Path(input_path)
        if p.exists():
            size = p.stat().st_size
            print(f"Veillez patienter...")
            print(f"🔄️")
            print(f"🔄️")
            print(f"🔄️")
            print(f"🔄️Converson du Fichier📄 {p.name} ({human_readable_size(size)}) en cours...")
            
    
    start_time = time.time()
    ok, msg = conv.convert_file_with_progress(input_path, out)
    elapsed = time.time() - start_time
    
    if ok:
        print(f"✅ {msg}")
        if not args.no_progress:
            print(f"⏱️  Temps de conversion: {elapsed / 60:.2f} min")
    else:
        print(f"❌ {msg}")

if __name__ == "__main__":
    main()